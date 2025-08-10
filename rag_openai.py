import os
import tempfile
import asyncio
import aiohttp
import numpy as np
from typing import List, Optional, Set, Dict, Callable, Any
import faiss
import requests
import json
from openai import AsyncOpenAI
from openai import RateLimitError, BadRequestError
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import logging
import fitz  # PyMuPDF for enhanced PDF processing
import re
from urllib.parse import urlparse, unquote
import pandas as pd
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv

try:
    from bs4 import BeautifulSoup  # HTML extraction
except Exception:
    BeautifulSoup = None

logger = logging.getLogger(__name__)
load_dotenv()

SEEN_URLS: Set[str] = set([
   "https://hackrx.blob.core.windows.net/hackrx/rounds/News.pdf?sv=2023-01-03&spr=https&st=2025-08-07T17%3A10%3A11Z&se=2026-08-08T17%3A10%3A00Z&sr=b&sp=r&sig=ybRsnfv%2B6VbxPz5xF7kLLjC4ehU0NF7KDkXua9ujSf0%3D",
   "https://hackrx.blob.core.windows.net/hackrx/rounds/FinalRound4SubmissionPDF.pdf?sv=2023-01-03&spr=https&st=2025-08-07T14%3A23%3A48Z&se=2027-08-08T14%3A23%3A00Z&sr=b&sp=r&sig=nMtZ2x9aBvz%2FPjRWboEOZIGB%2FaGfNf5TfBOrhGqSv4M%3D"
])

class RAGOpenAI:
    SYSTEM_PROMPT = (
        "You are a helpful assistant. First, use the provided context to understand the question. "
        "If the context provides a URL or mentions an action to get more information, use the "
        "available tools to perform that action. Then, answer the user's question based on all "
        "the information you have gathered. Be precise and concise."
    )

    def __init__(
        self,
        llm_model="gpt-5",
        embedding_model="text-embedding-3-small",
        top_k=3,
        collection_name="pdfs",
        backend: str = "faiss",
        faiss_use_gpu: bool = False,
        embed_batch_size: int = 128,
        use_enhanced_processing: bool = True,
        enhanced_chunk_size: int = 1000,
        enhanced_chunk_overlap: int = 300
    ):
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise RuntimeError("OPENAI_API_KEY not set")
        self.client = AsyncOpenAI(api_key=api_key)
        self.llm_model = llm_model
        self.embedding_model = embedding_model
        self.top_k = top_k
        self.collection_name = collection_name
        self.backend = backend
        self.faiss_use_gpu = bool(faiss_use_gpu)
        self.embed_batch_size = int(embed_batch_size)
        
        self.use_enhanced_processing = use_enhanced_processing
        self.enhanced_chunk_size = enhanced_chunk_size
        self.enhanced_chunk_overlap = enhanced_chunk_overlap

        self._faiss_index: Optional[faiss.Index] = None
        self._faiss_chunks: List[str] = []
        self._faiss_url: Optional[str] = None
        
        self.tools: List[Dict[str, Any]] = [
            {
                "type": "function",
                "function": {
                    "name": "make_http_request",
                    "description": "Makes an HTTP GET request to a specified URL to fetch real-time data like flight numbers or secret tokens. Only use for GET requests. i.e. https://register.hackrx.in/submissions/myFavouriteCity or ",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "url": {
                                "type": "string",
                                "description": "The URL to send the GET request to, e.g., 'https://register.hackrx.in/teams/public/flights/getThirdCityFlightNumber'. The details are in the documentation.",
                            },
                        },
                        "required": ["url"],
                    },
                },
            }
        ]
        self.available_tools: Dict[str, Callable[..., Any]] = {
            "make_http_request": self._make_http_request
        }

    async def _log(self, msg: str):
        logger.info(msg)

    async def _make_http_request(self, url: str) -> str:
        await self._log(f"Tool: Calling 'make_http_request' for URL: {url}")
        
        def sync_get():
            try:
                response = requests.get(url, timeout=15)
                response.raise_for_status()
                try:
                    return json.dumps(response.json())
                except json.JSONDecodeError:
                    return response.text
            except requests.exceptions.RequestException as e:
                return f"Error: The request failed. {e}"

        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(None, sync_get)
        await self._log(f"Tool: Request to {url} returned: {result[:200]}...")
        return result

    def _get_file_extension_from_url(self, url: str) -> str:
        parsed_url = urlparse(url)
        path = unquote(parsed_url.path)
        if '.' in path:
            ext = path.split('.')[-1].lower()
            # Treat unknown/asset-less extensions as HTML
            known = {'pdf','docx','doc','pptx','ppt','xlsx','xls','html','htm','txt','csv'}
            return ext if ext in known else 'html'
        # No extension -> assume HTML website
        return 'html'

    async def _download_file(self, url: str) -> tuple[str, str]:
        if os.path.isfile(url):
            await self._log(f"Local file found: {url}")
            ext = url.split('.')[-1].lower() if '.' in url else 'pdf'
            return url, ext
        
        file_ext = self._get_file_extension_from_url(url)
        fd, tmp = tempfile.mkstemp(suffix=f".{file_ext}")
        os.close(fd)
        
        await self._log(f"Downloading: {url}")
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=60) as r:
                r.raise_for_status()
                with open(tmp, "wb") as f:
                    f.write(await r.read())
        await self._log(f"Downloaded to: {tmp} (type: {file_ext})")
        return tmp, file_ext

    def _clean_text(self, text: str) -> str:
        if not text: return ""
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n\d+\n', '\n', text)
        text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
        text = re.sub(r'\n+', '\n', text)
        return text.strip()

    async def _extract_text_from_html(self, file_path: str) -> str:
        """Extract readable text from an HTML file downloaded to disk."""
        def extract():
            try:
                with open(file_path, "rb") as f:
                    raw = f.read()
            except Exception as e:
                logger.warning(f"Failed to read HTML file: {e}")
                return ""
            html = raw.decode("utf-8", errors="ignore")
            if BeautifulSoup:
                soup = BeautifulSoup(html, "html.parser")
                for tag in soup(["script", "style", "noscript"]):
                    tag.decompose()
                title = soup.title.string.strip() if soup.title and soup.title.string else ""
                text = soup.get_text(separator="\n")
                combined = f"{title}\n\n{text}" if title else text
                return self._clean_text(combined)
            # Fallback: naive strip tags
            text = re.sub(r"(?is)<(script|style).*?>.*?(</\1>)", " ", html)
            text = re.sub(r"<[^>]+>", " ", text)
            return self._clean_text(text)

        return await asyncio.get_event_loop().run_in_executor(None, extract)

    async def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        def extract():
            doc = Document(file_path)
            full_text = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            
            # Extract tables
            for table in doc.tables:
                table_text = "\n\n**Table:**\n"
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_text += f"| {row_text} |\n"
                full_text.append(table_text)
            
            return self._clean_text("\n\n".join(full_text))
        
        return await asyncio.get_event_loop().run_in_executor(None, extract)

    async def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        def extract():
            prs = Presentation(file_path)
            full_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n\n**Slide {slide_num}:**\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += f"{shape.text.strip()}\n"
                    
                    # Extract table data if present
                    if shape.has_table:
                        table_text = "\n**Table in slide:**\n"
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                table_text += f"| {row_text} |\n"
                        slide_text += table_text
                
                if slide_text.strip() != f"**Slide {slide_num}:**":
                    full_text.append(slide_text)
            
            return self._clean_text("\n".join(full_text))
        
        return await asyncio.get_event_loop().run_in_executor(None, extract)

    async def _extract_text_from_xlsx(self, file_path: str) -> str:
        """Extract text from Excel file."""
        def extract():
            full_text = []
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                sheet_text = f"\n\n**Sheet: {sheet_name}**\n"
                
                # Add column headers
                if not df.empty:
                    headers = " | ".join(str(col) for col in df.columns)
                    sheet_text += f"| {headers} |\n"
                    sheet_text += f"|{' --- |' * len(df.columns)}\n"
                    
                    # Add rows (limit to avoid huge chunks)
                    for idx, row in df.head(100).iterrows():  # Limit to first 100 rows
                        row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row)
                        sheet_text += f"| {row_text} |\n"
                    
                    if len(df) > 100:
                        sheet_text += f"\n... and {len(df) - 100} more rows\n"
                
                full_text.append(sheet_text)
                if not sheet_text.strip():
                    logger.warning(f"No usable text extracted from {file_path}")
                    return ["No content"]

            
            return self._clean_text("\n".join(full_text))
        
        return await asyncio.get_event_loop().run_in_executor(None, extract)

    async def _extract_text_from_pdf(self, file_path: str) -> str:
        """Enhanced PDF text extraction using PyMuPDF."""
        def extract():
            doc = fitz.open(file_path)
            full_text = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text with layout preservation
                text = page.get_text("text")
                
                # Extract table data
                tables = page.find_tables()
                for table in tables:
                    table_data = table.extract()
                    if table_data:
                        # Convert table to markdown format
                        table_text = "\n\n**Table:**\n"
                        for row in table_data:
                            if row and any(cell and cell.strip() for cell in row):
                                table_text += "| " + " | ".join(str(cell or "") for cell in row) + " |\n"
                        text += table_text
                
                # Get image information (for context)
                images = page.get_images()
                if images:
                    text += f"\n\n[Page contains {len(images)} image(s)]"
                
                full_text.append(self._clean_text(text))
            
            doc.close()
            return "\n\n".join(full_text)
        
        return await asyncio.get_event_loop().run_in_executor(None, extract)

    async def _extract_text_from_file(self, file_path: str, file_ext: str) -> str:
        try:
            if file_ext == 'pdf': return await self._extract_text_from_pdf(file_path)
            elif file_ext in ['docx', 'doc']: return await self._extract_text_from_docx(file_path)
            elif file_ext in ['pptx', 'ppt']: return await self._extract_text_from_pptx(file_path)
            elif file_ext in ['xlsx', 'xls']: return await self._extract_text_from_xlsx(file_path)
            elif file_ext in ['html', 'htm']: return await self._extract_text_from_html(file_path)
            else:
                await self._log(f"Unsupported file type: {file_ext}")
                return f"Unsupported file type: {file_ext}"
        except Exception as e:
            await self._log(f"Error extracting text from {file_ext}: {e}")
            return f"Error extracting content from file: {str(e)}"

    async def _chunk_texts(self, file_path: str, file_ext: str, use_enhanced: bool) -> List[str]:
        # For PDFs, we keep a fast loader path when not enhanced
        if file_ext == 'pdf' and not use_enhanced:
            docs = PyPDFLoader(file_path).load()
            splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=80)
            return [d.page_content for d in splitter.split_documents(docs)]
        # For all other types (incl. HTML) or enhanced PDFs, extract then split
        text_content = await self._extract_text_from_file(file_path, file_ext)
        if not text_content or not isinstance(text_content, str) or not text_content.strip():
            return [f"No readable content in {file_ext.upper()}."]
        chunk_size = self.enhanced_chunk_size if use_enhanced and self.use_enhanced_processing else 800
        chunk_overlap = self.enhanced_chunk_overlap if use_enhanced and self.use_enhanced_processing else 120
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        return splitter.split_text(text_content)

    async def _embed_texts_async(self, texts: List[str]) -> np.ndarray:
        if not texts: return np.array([], dtype=np.float32).reshape(0, 1536) # Reshape to handle empty case
        vecs = []
        for i in range(0, len(texts), self.embed_batch_size):
            batch = texts[i:i + self.embed_batch_size]
            resp = await self.client.embeddings.create(model=self.embedding_model, input=batch)
            vecs.extend([np.array(d.embedding, dtype=np.float32) for d in resp.data])
        
        def process_embeddings():
            mat = np.vstack(vecs)
            faiss.normalize_L2(mat)
            return mat
        return await asyncio.get_event_loop().run_in_executor(None, process_embeddings)

    async def _build_faiss_index(self, chunks: List[str]):
        emb = await self._embed_texts_async(chunks)
        if emb.shape[0] == 0:
            self._faiss_index = None
            self._faiss_chunks = []
            await self._log("No embeddings generated; FAISS index is empty.")
            return

        d = emb.shape[1]
        index = faiss.IndexFlatIP(d)
        index.add(emb)
        self._faiss_index = index
        self._faiss_chunks = chunks
        await self._log("FAISS index built on CPU.")
    
    async def ensure_vector_store_for_url(self, url: str, seen_urls: Set[str] = None):
        if self.backend == "faiss" and self._faiss_url == url: return
        await self._log(f"Ensuring vector store for: {url}")
        is_seen_url = seen_urls is not None and url in seen_urls
        path, file_ext = await self._download_file(url)
        try:
            chunks = await self._chunk_texts(path, file_ext, use_enhanced=is_seen_url)
            await self._build_faiss_index(chunks)
            self._faiss_url = url
        finally:
            if not os.path.isfile(url) and os.path.exists(path): os.remove(path)

    async def _answer_one_question_with_tools(self, question: str, context: str) -> str:
        messages: List[Dict[str, Any]] = [
            {"role": "system", "content": self.SYSTEM_PROMPT},
            {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {question}"}
        ]
        
        response = await self.client.chat.completions.create(
            model=self.llm_model,
            messages=messages,
            tools=self.tools,
            tool_choice="auto",
        )
        
        response_message = response.choices[0].message
        tool_calls = response_message.tool_calls

        if not tool_calls:
            return response_message.content.strip() if response_message.content else "No answer found."

        messages.append(response_message)
        
        for tool_call in tool_calls:
            function_name = tool_call.function.name
            function_to_call = self.available_tools.get(function_name)
            if function_to_call:
                function_args = json.loads(tool_call.function.arguments)
                function_response = await function_to_call(**function_args)
                messages.append({
                    "tool_call_id": tool_call.id,
                    "role": "tool",
                    "name": function_name,
                    "content": function_response,
                })
            else:
                await self._log(f"Warning: LLM requested unknown tool '{function_name}'")

        second_response = await self.client.chat.completions.create(
            model=self.llm_model, messages=messages
        )
        return second_response.choices[0].message.content.strip()

    async def answer_questions(self, questions: List[str], scope_url: Optional[str] = None) -> List[str]:
        if self.backend == "faiss":
            # Fallback to no-context chat completion if index is missing/empty
            if self._faiss_index is None:
                await self._log("FAISS index not ready; answering without context.")
                tasks = [self._answer_one_question_with_tools(q, "") for q in questions]
                return await asyncio.gather(*tasks)

            q_emb = await self._embed_texts_async(questions)
            tasks = []
            for i, question_embedding in enumerate(q_emb):
                async def create_task(q_idx, q_vec):
                    scores, idxs = await asyncio.get_event_loop().run_in_executor(
                        None, lambda: self._faiss_index.search(q_vec.reshape(1, -1), self.top_k)
                    )
                    context = "\n\n".join(self._faiss_chunks[j] for j in idxs[0] if 0 <= j < len(self._faiss_chunks))
                    return await self._answer_one_question_with_tools(questions[q_idx], context)
                tasks.append(create_task(i, question_embedding))
            return await asyncio.gather(*tasks)

        return ["Remote backend not supported with tools in this version." for _ in questions]

    async def process_pdf_and_answer(self, pdf_url: str, questions: List[str], use_memory: bool = True, seen_urls: Set[str] = None) -> List[str]:
        await self.ensure_vector_store_for_url(pdf_url, seen_urls=seen_urls)
        return await self.answer_questions(questions, scope_url=pdf_url)
